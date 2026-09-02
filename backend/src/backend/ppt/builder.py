"""Deterministic slide-spec -> .pptx renderer.

Pure and synchronous: no disk, no network, no LLM. Safe to call inside a
request handler (~50-200 ms for a full deck). Starts from python-pptx's
built-in template and applies a light Medha accent; a designed template
`.pptx` can be dropped in later behind this same function.
"""

import io
import logging
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from backend.ppt.schema import DeckSpec, SlideSpec

logger = logging.getLogger("backend.ppt")

RENDER_VERSION = "render-v1"

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

# built-in template layout indices: 0 = Title Slide, 1 = Title and Content
_TITLE_LAYOUT = 0
_CONTENT_LAYOUT = 1

_ACCENT_RGB = RGBColor(0x8A, 0x3D, 0x2B)  # deep terracotta
_FOOTER_RGB = RGBColor(0x8A, 0x8A, 0x8A)


def slugify_filename(text: str) -> str:
    """ASCII-safe basename for a downloaded deck (no extension)."""
    slug = re.sub(r"[^A-Za-z0-9]+", "-", text or "").strip("-").lower()
    return slug or "medha-slides"


def _style_title(shape) -> None:
    if shape is None or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = _ACCENT_RGB
            run.font.bold = True


def _body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def _add_footer(slide, prs, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4),
        prs.slide_height - Inches(0.5),
        prs.slide_width - Inches(0.8),
        Inches(0.35),
    )
    box.text_frame.text = text
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(9)
    run.font.color.rgb = _FOOTER_RGB


def _content_slide(prs, spec: SlideSpec, footer: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_CONTENT_LAYOUT])
    slide.shapes.title.text = spec.heading
    _style_title(slide.shapes.title)

    body = _body_placeholder(slide)
    if body is not None and body.has_text_frame and spec.bullets:
        tf = body.text_frame
        tf.text = spec.bullets[0]
        for bullet in spec.bullets[1:]:
            tf.add_paragraph().text = bullet

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if footer:
        _add_footer(slide, prs, footer)


def render_pptx(deck: DeckSpec, *, footer: str | None = None) -> bytes:
    """Render a validated `DeckSpec` to `.pptx` bytes."""
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    title_slide.shapes.title.text = deck.title
    _style_title(title_slide.shapes.title)
    subtitle_ph = _body_placeholder(title_slide)
    if subtitle_ph is not None:
        subtitle_ph.text = deck.subtitle

    for spec in deck.slides:
        _content_slide(prs, spec, footer)

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    logger.info(
        "ppt render slides=%d bytes=%d version=%s",
        len(deck.slides),
        len(data),
        RENDER_VERSION,
    )
    return data

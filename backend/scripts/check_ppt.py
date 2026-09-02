"""Standalone sanity check for the slide-spec parser + renderer.

This repo has no pytest setup, so this is a plain runnable script:

    uv run python scripts/check_ppt.py

Exits non-zero on the first failed assertion.
"""

import io

from pptx import Presentation

from backend.ppt.builder import render_pptx, slugify_filename
from backend.ppt.schema import (
    MAX_BULLETS_PER_SLIDE,
    MAX_SLIDES,
    DeckParseError,
    parse_deck,
)

_SPEC = {
    "_prompt_version": "ppt-v1",
    "title": "Friction",
    "subtitle": "Class 8 - Science",
    "slides": [
        {
            "layout": "bullets",
            "heading": "What is friction?",
            "bullets": ["Opposes motion", "Acts along the surface"],
            "notes": "Slide a book and let it stop.",
        },
        {
            "layout": "weird-unknown",  # -> falls back to bullets
            "heading": "Types",
            "bullets": ["Static", "Sliding", "Rolling"],
            "notes": "Rolling is smallest.",
        },
        {"heading": "", "bullets": []},  # -> dropped
    ],
}


def check_parse_ok() -> None:
    deck = parse_deck(_SPEC)
    assert deck.title == "Friction", deck.title
    assert len(deck.slides) == 2, "empty slide should have been dropped"
    assert deck.slides[1].layout == "bullets", "unknown layout should fall back"


def check_parse_clamps() -> None:
    big = {
        "title": "x" * 500,
        "slides": [{"heading": "h", "bullets": ["b" * 400] + ["x"] * 50}] * 50,
    }
    deck = parse_deck(big)
    assert len(deck.title) == 160
    assert len(deck.slides) == MAX_SLIDES
    assert len(deck.slides[0].bullets) == MAX_BULLETS_PER_SLIDE
    assert len(deck.slides[0].bullets[0]) == 200


def check_parse_rejects() -> None:
    for bad in [{}, {"slides": []}, {"slides": [{"heading": "", "bullets": []}]}, "nope"]:
        try:
            parse_deck(bad)
        except DeckParseError:
            continue
        raise AssertionError(f"expected DeckParseError for {bad!r}")


def check_render() -> None:
    deck = parse_deck(_SPEC)
    data = render_pptx(deck, footer="Medha - Class 8 Science")
    assert data[:2] == b"PK", "not a zip / pptx"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == len(deck.slides) + 1, "expected 1 title + content slides"
    assert prs.slides[0].shapes.title.text == "Friction"
    assert prs.slides[1].shapes.title.text == "What is friction?"
    assert prs.slides[1].has_notes_slide
    assert prs.slides[1].notes_slide.notes_text_frame.text.strip()


def check_slug() -> None:
    assert slugify_filename("प्रकाश  संश्लेषण!!") == "medha-slides"
    assert slugify_filename("How Plants Make Their Food") == "how-plants-make-their-food"
    assert slugify_filename("") == "medha-slides"


def main() -> None:
    for fn in (
        check_parse_ok,
        check_parse_clamps,
        check_parse_rejects,
        check_render,
        check_slug,
    ):
        fn()
        print(f"ok  {fn.__name__}")
    print("\nall checks passed")


if __name__ == "__main__":
    main()
